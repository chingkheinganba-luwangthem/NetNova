# -*- coding: utf-8 -*-
"""
PPTX Redesign - v6
Matches the reference slide 1 image pixel-accurately:
  - Soft blue gradient background (top-to-bottom)
  - Full-width dark navy rounded-corner title banner
  - Thin gold accent bar just below title
  - Gold border + rounded corners + shadow on ALL images
  - Bullet-point ► formatting preserved (only font/color changed)
  - Consistent Times New Roman, proper line spacing, padding
  - Modern table: navy header, alternating rows, clean borders
  - THANK YOU: solid navy bg, large white centred text
"""

from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from lxml import etree

INPUT  = r"d:\Globizs\PDF\Legal Metrology-Presentation-Fnl.pptx"
OUTPUT = r"d:\Globizs\PDF\Legal Metrology-Redesigned-v6.pptx"

# ── PALETTE ───────────────────────────────────────────────────────────────────
BG_TOP      = "C8D8EE"   # gradient top   – muted periwinkle
BG_BOT      = "EEF4FC"   # gradient bottom – near-white blue
NAVY        = RGBColor(0x16, 0x25, 0x4D)
GOLD        = RGBColor(0xC9, 0xA0, 0x48)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
BODY_FG     = RGBColor(0x16, 0x25, 0x4D)
CAPTION_FG  = RGBColor(0x16, 0x25, 0x4D)
TABLE_HDR   = RGBColor(0x16, 0x25, 0x4D)
TABLE_HDR_FG= RGBColor(0xFF, 0xFF, 0xFF)
TABLE_EVEN  = RGBColor(0xEE, 0xF2, 0xFF)
TABLE_ODD   = RGBColor(0xFF, 0xFF, 0xFF)
TABLE_SUB   = RGBColor(0xC7, 0xD2, 0xFE)
TABLE_TOT   = RGBColor(0xE0, 0xE7, 0xFF)
TABLE_BRD   = RGBColor(0xA5, 0xB4, 0xFC)

FONT          = "Times New Roman"
SZ_TITLE      = 26   # pt
SZ_BODY       = 13   # pt
SZ_CAPTION    = 11   # pt
SZ_TABLE      = 10   # pt
LS_TITLE      = 28   # line spacing pt
LS_BODY       = 20
LS_CAPTION    = 15
LS_TABLE      = 13
SPC_BODY      = 4    # paragraph spc before/after
TF_PAD        = 182880   # 0.2 inch in EMU

# ── HELPERS ───────────────────────────────────────────────────────────────────

def rgb_hex(c): return "{:02X}{:02X}{:02X}".format(c[0], c[1], c[2])

# ── GRADIENT BACKGROUND ───────────────────────────────────────────────────────

def set_gradient_bg(slide, top_hex=BG_TOP, bot_hex=BG_BOT):
    """
    Apply a top-to-bottom linear gradient to the slide background.
    Matches the reference image's light blue-to-white gradient.
    """
    cSld = slide._element.find(qn('p:cSld'))
    if cSld is None:
        return
    bg = cSld.find(qn('p:bg'))
    if bg is None:
        bg = etree.SubElement(cSld, qn('p:bg'))
    bgPr = bg.find(qn('p:bgPr'))
    if bgPr is None:
        bgPr = etree.SubElement(bg, qn('p:bgPr'))
    # clear existing fill nodes
    for tag in [qn('a:solidFill'), qn('a:gradFill'), qn('a:noFill'),
                qn('a:pattFill'), qn('a:blipFill')]:
        for old in bgPr.findall(tag):
            bgPr.remove(old)
    gradFill = etree.SubElement(bgPr, qn('a:gradFill'))
    gsLst    = etree.SubElement(gradFill, qn('a:gsLst'))
    for pos, hexval in [('0', top_hex), ('100000', bot_hex)]:
        gs = etree.SubElement(gsLst, qn('a:gs'))
        gs.set('pos', pos)
        clr = etree.SubElement(gs, qn('a:srgbClr'))
        clr.set('val', hexval)
    lin = etree.SubElement(gradFill, qn('a:lin'))
    lin.set('ang', '5400000')   # 90°: top → bottom
    lin.set('scaled', '0')
    # ensure bgRef is removed (it overrides bgPr otherwise)
    for old in bg.findall(qn('p:bgRef')):
        bg.remove(old)


# ── LINE SPACING / PARAGRAPH ───────────────────────────────────────────────────

def set_line_spacing(para, pt):
    pPr = para._p.get_or_add_pPr()
    for old in pPr.findall(qn('a:lnSpc')):
        pPr.remove(old)
    lnSpc  = etree.SubElement(pPr, qn('a:lnSpc'))
    spcPts = etree.SubElement(lnSpc, qn('a:spcPts'))
    spcPts.set('val', str(int(pt * 100)))

def set_para_spc(para, before_pt, after_pt):
    pPr = para._p.get_or_add_pPr()
    pPr.set('spcBef', str(int(before_pt * 100)))
    pPr.set('spcAft', str(int(after_pt  * 100)))


# ── TEXT FRAME STYLING ────────────────────────────────────────────────────────

def style_tf(tf, color, bold=False, size_pt=SZ_BODY,
             line_pt=LS_BODY, spc=SPC_BODY, align=None, pad=True):
    """
    Style every run in a text frame.
    Does NOT touch bullet XML — preserves ► arrows and other list markers.
    """
    if pad:
        tf.margin_left   = Emu(TF_PAD)
        tf.margin_right  = Emu(TF_PAD)
        tf.margin_top    = Emu(TF_PAD // 2)
        tf.margin_bottom = Emu(TF_PAD // 2)
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    for para in tf.paragraphs:
        set_line_spacing(para, line_pt)
        set_para_spc(para, spc, spc)
        if align is not None:
            para.alignment = align
        for run in para.runs:
            run.font.name      = FONT
            run.font.color.rgb = color
            run.font.bold      = bold
            run.font.size      = Pt(size_pt)


# ── GEOMETRY / SHADOW / BORDER ────────────────────────────────────────────────

def _spPr(shape):
    sp = shape._element
    r  = sp.find(qn('p:spPr'))
    if r is None:
        r = sp.find('.//{http://schemas.openxmlformats.org/drawingml/2006/picture}spPr')
    return r

def apply_round_rect(spPr, adj_pct=10):
    for tag in [qn('a:prstGeom'), qn('a:custGeom')]:
        for old in spPr.findall(tag):
            spPr.remove(old)
    pg = etree.Element(qn('a:prstGeom'))
    pg.set('prst', 'roundRect')
    av = etree.SubElement(pg, qn('a:avLst'))
    gd = etree.SubElement(av, qn('a:gd'))
    gd.set('name', 'adj')
    gd.set('fmla', f'val {adj_pct * 1000}')
    before = next((c for c in spPr
                   if c.tag in [qn('a:ln'), qn('a:effectLst')]), None)
    if before is not None:
        spPr.insert(list(spPr).index(before), pg)
    else:
        spPr.append(pg)

def apply_shadow(spPr):
    for old in spPr.findall(qn('a:effectLst')):
        spPr.remove(old)
    eLst   = etree.SubElement(spPr, qn('a:effectLst'))
    shadow = etree.SubElement(eLst, qn('a:outerShdw'))
    shadow.set('blurRad', '76200')
    shadow.set('dist',    '45720')
    shadow.set('dir',     '5400000')
    shadow.set('algn',    'tl')
    shadow.set('rotWithShape', '0')
    sc = etree.SubElement(shadow, qn('a:srgbClr'))
    sc.set('val', '162540')
    al = etree.SubElement(sc, qn('a:alpha'))
    al.set('val', '30000')

def apply_gold_border(spPr, width_pt=3):
    """Add a visible gold border around the image (matches reference frame)."""
    for old in spPr.findall(qn('a:ln')):
        spPr.remove(old)
    ln = etree.SubElement(spPr, qn('a:ln'))
    ln.set('w', str(int(width_pt * 12700)))  # pt → EMU
    ln.set('cap', 'flat')
    ln.set('cmpd', 'sng')
    sf = etree.SubElement(ln, qn('a:solidFill'))
    sc = etree.SubElement(sf, qn('a:srgbClr'))
    sc.set('val', rgb_hex(GOLD))

def style_image(shape, adj_pct=10):
    """Rounded corners + gold border + soft shadow on any image shape."""
    sp = _spPr(shape)
    if sp is None:
        return
    apply_round_rect(sp, adj_pct)
    apply_gold_border(sp, width_pt=3)
    apply_shadow(sp)


# ── TITLE BANNER ──────────────────────────────────────────────────────────────

def style_title_banner(shape):
    shape.fill.solid()
    shape.fill.fore_color.rgb = NAVY
    shape.line.fill.background()
    sp = _spPr(shape)
    if sp is not None:
        apply_round_rect(sp, adj_pct=5)
    style_tf(shape.text_frame,
             color=WHITE, bold=True, size_pt=SZ_TITLE,
             line_pt=LS_TITLE, spc=4, align=PP_ALIGN.CENTER, pad=True)

def add_gold_bar(slide, title_shape):
    bar = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        title_shape.left,
        title_shape.top + title_shape.height,
        title_shape.width,
        Emu(50800)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = GOLD
    bar.line.fill.background()


# ── TABLE ─────────────────────────────────────────────────────────────────────

def style_table(shape):
    if not shape.has_table:
        return
    table = shape.table
    nrows = len(table.rows)
    for r, row in enumerate(table.rows):
        first    = row.cells[0].text.strip().upper()
        is_sub   = (len(first) == 1 and first in ['A','B','C'])
        is_total = ('TOTAL' in first or r == nrows - 1)
        for c, cell in enumerate(row.cells):
            tc   = cell._tc
            tcPr = tc.find(qn('a:tcPr'))
            if tcPr is None:
                tcPr = etree.SubElement(tc, qn('a:tcPr'))
            tcPr.set('anchor', 'ctr')
            for tag in [qn('a:solidFill'), qn('a:noFill'), qn('a:gradFill')]:
                for old in tcPr.findall(tag):
                    tcPr.remove(old)
            if r == 0:
                bg, fg, bold = rgb_hex(TABLE_HDR), TABLE_HDR_FG, True
            elif is_sub:
                bg, fg, bold = rgb_hex(TABLE_SUB), NAVY, True
            elif is_total:
                bg, fg, bold = rgb_hex(TABLE_TOT), NAVY, True
            elif r % 2 == 0:
                bg, fg, bold = rgb_hex(TABLE_EVEN), BODY_FG, False
            else:
                bg, fg, bold = rgb_hex(TABLE_ODD),  BODY_FG, False
            sf = etree.SubElement(tcPr, qn('a:solidFill'))
            sc = etree.SubElement(sf, qn('a:srgbClr'))
            sc.set('val', bg)
            for btag in [qn('a:lnL'),qn('a:lnR'),qn('a:lnT'),qn('a:lnB')]:
                for old in tcPr.findall(btag):
                    tcPr.remove(old)
                ln = etree.SubElement(tcPr, btag)
                ln.set('w', '12700'); ln.set('cap','flat'); ln.set('cmpd','sng')
                bsf = etree.SubElement(ln, qn('a:solidFill'))
                bsc = etree.SubElement(bsf, qn('a:srgbClr'))
                bsc.set('val', rgb_hex(TABLE_BRD))
            for attr, val in [('marL','91440'),('marR','91440'),
                              ('marT','45720'),('marB','45720')]:
                tcPr.set(attr, val)
            for para in cell.text_frame.paragraphs:
                set_line_spacing(para, LS_TABLE)
                set_para_spc(para, 2, 2)
                for run in para.runs:
                    run.font.name      = FONT
                    run.font.color.rgb = fg
                    run.font.bold      = bold
                    run.font.size      = Pt(SZ_TABLE)


# ── TITLE DETECTION ───────────────────────────────────────────────────────────

def is_title(shape):
    if shape.shape_type != 14:
        return False
    try:
        ph = shape.placeholder_format
        return ph is not None and ph.idx == 0
    except Exception:
        return False


# ── MAIN ─────────────────────────────────────────────────────────────────────

def redesign_v6(input_path, output_path):
    prs = Presentation(input_path)

    for slide_idx, slide in enumerate(prs.slides):
        # ── BACKGROUND ───────────────────────────────────────────────────────
        set_gradient_bg(slide)   # gradient, no ovals

        # ── THANK YOU SLIDE ──────────────────────────────────────────────────
        if slide_idx == 6:
            set_gradient_bg(slide, top_hex='16254D', bot_hex='1E3A6E')
            for shape in slide.shapes:
                if shape.has_text_frame and shape.text.strip():
                    try:
                        shape.fill.background()
                    except Exception:
                        pass
                    style_tf(shape.text_frame,
                             color=WHITE, bold=True, size_pt=60,
                             line_pt=70, spc=0, align=PP_ALIGN.CENTER, pad=False)
            continue

        title_shape = None

        for shape in slide.shapes:
            stype = shape.shape_type

            if is_title(shape):
                title_shape = shape
                style_title_banner(shape)

            elif stype == 13:                              # picture
                style_image(shape)

            elif stype == 14 and not shape.has_text_frame: # picture placeholder
                style_image(shape)

            elif stype == 14 and shape.has_text_frame:     # body placeholder
                try:
                    shape.fill.background()
                except Exception:
                    pass
                style_tf(shape.text_frame,
                         color=BODY_FG, size_pt=SZ_BODY, line_pt=LS_BODY, spc=SPC_BODY)

            elif stype == 17 and shape.has_text_frame:     # text boxes / captions
                try:
                    shape.fill.background()
                except Exception:
                    pass
                txt = shape.text.strip()
                if len(txt) < 120:
                    style_tf(shape.text_frame,
                             color=CAPTION_FG, bold=True,
                             size_pt=SZ_CAPTION, line_pt=LS_CAPTION, spc=2,
                             align=PP_ALIGN.CENTER)
                else:
                    style_tf(shape.text_frame,
                             color=BODY_FG, size_pt=SZ_BODY, line_pt=LS_BODY, spc=SPC_BODY)

            elif stype == 19:                              # table
                style_table(shape)

        if title_shape is not None:
            add_gold_bar(slide, title_shape)

    prs.save(output_path)
    print("Saved: " + output_path)


if __name__ == "__main__":
    redesign_v6(INPUT, OUTPUT)
