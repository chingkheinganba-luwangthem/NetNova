import sys
import json
import os
from pptx import Presentation

def extract(pptx_path, out_dir):
    prs = Presentation(pptx_path)
    slides_data = []
    
    os.makedirs(out_dir, exist_ok=True)
    
    for i, slide in enumerate(prs.slides):
        slide_info = {
            "slide_number": i + 1,
            "texts": [],
            "images": []
        }
        
        image_count = 0
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_info["texts"].append(shape.text.strip())
            
            if shape.shape_type == 13: # Picture
                image_count += 1
                image = shape.image
                image_bytes = image.blob
                image_ext = image.ext
                image_filename = f"slide_{i+1}_image_{image_count}.{image_ext}"
                image_path = os.path.join(out_dir, image_filename)
                with open(image_path, "wb") as f:
                    f.write(image_bytes)
                slide_info["images"].append(image_filename)
                
            if shape.shape_type == 19: # Table
                table_data = []
                for row in shape.table.rows:
                    row_data = []
                    for cell in row.cells:
                        row_data.append(cell.text.strip() if cell.text else "")
                    table_data.append(row_data)
                if "tables" not in slide_info:
                    slide_info["tables"] = []
                slide_info["tables"].append(table_data)
                
        slides_data.append(slide_info)
        
    with open(os.path.join(out_dir, "presentation_data.json"), "w", encoding="utf-8") as f:
        json.dump(slides_data, f, indent=2, ensure_ascii=False)
        
    print(f"Extracted {len(slides_data)} slides. Data saved to {out_dir}")

if __name__ == "__main__":
    pptx_path = r"d:\Globizs\PDF\Legal Metrology-Presentation-Fnl.pptx"
    out_dir = r"d:\Projects\NetNova\scratch\pptx_data"
    extract(pptx_path, out_dir)
